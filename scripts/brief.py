"""Generate a slide deck and an HTML brief from an insights report - no number invented.

Why this exists
---------------
The last mile of "retrieve insights" is a document someone presents. That is
where numbers get retyped, rounded, reordered and occasionally made up, and it
is the one artifact a buyer actually holds. So the deck is GENERATED from the
insights JSON, the report's figures were compiled and conserved by the
ontology, and the chain from a slide back to an executed query is unbroken and
machine-checkable.

What is promised, precisely
---------------------------
* Every value on the deck and in the HTML is rendered by ONE formatter, at 2 dp
  with thousands separators and the unit appended. The report is validated up
  front to hold its values at 2 dp already, so rendering at 2 dp is verbatim
  and nothing is re-rounded - not on the deck, not in the HTML.
* A table shows all of an insight's rows, or at most MAX_ROWS of them with an
  explicit line saying how many of the groups are shown and what the rest of
  the total amounts to. That remainder is derived - total minus the shown
  rows, rounded to 2 dp - and is the only figure here that is not a report
  field verbatim; it is computed in one place and named for what it is.
* The footer repeats the share percentage exactly as the headline states it -
  one string, not a second rendering of top_share that can land on the other
  side of a rounding boundary.
* Every field interpolated into HTML passes through one escaping helper.
* A report that is missing a field, holds an empty headline, a None figure or
  an insight with no rows is refused up front with "FAIL <insight id>: <what is
  wrong>" - never a traceback.
* After writing the deck, main() reads it back with python-pptx and checks that
  every numeric token on every slide is one the report accounts for (row
  values, totals, group counts, the headline's share, the truncation line, the
  ontology counts, and digits inside labels the report itself carries). PASS
  is printed only after that read-back is clean. The test calls the same
  verify_deck(). With --html only there is no deck to read back, so the run
  says "wrote HTML; no deck built" and never PASS.

What it deliberately does not do
--------------------------------
It does not phrase. Headlines are the insight templates' own sentences, which
are factual and a little wooden. A model may later rewrite them for a reader;
it will do so with the figures pinned, because a rewrite that changes a number
is a wrong number, and the read-back would catch it.

  python scripts/brief.py --insights data/lake/_reports/insights_aw.json \\
      --pptx data/lake/_reports/aw_brief.pptx --html data/lake/_reports/aw_brief.html

Exit 0 with PASS when the deck was written and read back clean.
"""

from __future__ import annotations

import argparse
import html
import json
import math
import re
import sys
from pathlib import Path
from typing import Any

MAX_ROWS = 8
MAX_REFUSALS = 8

# Any numeric token: integers, 1 dp, 2 dp or more, with or without thousands
# separators. Digits glued to letters (AdventureWorks2025, AddressLine1) are
# names, not figures, and are not matched; the same function builds both the
# allowed set and the seen set, so the two sides agree on what a number is.
_NUMBER = re.compile(r"(?<![A-Za-z_\d.,])\d[\d,]*(?:\.\d+)?(?![A-Za-z_\d])")
# The share as the insight templates print it: "... carry 86.2 pct of ...",
# "... alone is 45.0 pct of ...", "9.1 pct of ... lands on ...".
_PCT_IN_HEADLINE = re.compile(r"(?<![\d.])(\d+(?:\.\d+)?) pct\b")
_SHARE_LABEL = {"concentration": "top share", "dominance": "largest group",
                "unknown_bucket": "NULL share"}


# --------------------------------------------------------------------------
# one rendering of every figure, used by the deck, the HTML and the check
# --------------------------------------------------------------------------


def _num(v: float) -> str:
    """The one number formatter: 2 dp, thousands separators. Verbatim for a 2 dp value."""
    return f"{v:,.2f}"


def _fmt(v: float, unit: str) -> str:
    """A value with its unit appended: '43,909,437.51 USD'."""
    return f"{_num(v)} {unit}".strip()


def _esc(x: Any) -> str:
    """Every field interpolated into HTML goes through here - labels, counts, names, SQL."""
    return html.escape(str(x), quote=True)


def _numbers(text: str) -> set[str]:
    """Every numeric token in text, thousands separators removed."""
    return {m.replace(",", "") for m in _NUMBER.findall(text)}


def _shown(i: dict[str, Any]) -> list[list[Any]]:
    return list(i["rows"][:MAX_ROWS])


def _share_text(i: dict[str, Any]) -> str:
    """The share as the headline states it - one string, never re-derived.

    insights.py prints the headline's percentage from the unrounded share and
    stores top_share rounded to 4 dp; rendering the field again can land on
    the other side of a rounding boundary (50.1 in the headline, 50.0 in the
    footer). The footer therefore repeats the headline's own token, labelled
    by kind. A headline without a percentage gets top_share rendered once,
    with nothing to disagree with.
    """
    found = _PCT_IN_HEADLINE.findall(str(i["headline"]))
    if found:
        # the last token: concentration and dominance headlines put labels
        # (which may themselves say "pct") before the share
        return f"{_SHARE_LABEL.get(i.get('kind'), 'share')} {found[-1]} pct"
    return f"top share {i['top_share'] * 100:.1f} pct"


def _remainder_text(i: dict[str, Any]) -> str:
    """What the table does not show, so the rows and the footer total reconcile.

    rows are the report's top labelled groups (at most MAX_ROWS are shown),
    n_groups counts the labelled groups, and total is the measure over every
    fact row, including rows with no value for the dimension. The remainder is
    total minus the shown rows, rounded to 2 dp: a derived figure, stated as
    one. When every group is shown and the remainder is within the 2 dp
    rounding noise of n+1 figures, nothing is stated - 0.01 would be an
    artifact, not an amount.
    """
    shown = _shown(i)
    rest = int(i["n_groups"]) - len(shown)
    rem = round(float(i["total"]) - math.fsum(float(v) for _, v in shown), 2)
    if rem == 0:
        rem = 0.0  # never '-0.00'
    if rest > 0:
        groups = "group" if rest == 1 else "groups"
        return (f"showing {len(shown)} of {i['n_groups']} groups; the other {rest} {groups} "
                f"and any rows with no {i['dimension']} value sum to {_fmt(rem, i['unit'])}")
    if abs(rem) > 0.005 * (len(shown) + 1):
        return (f"all {i['n_groups']} groups shown; rows with no {i['dimension']} value "
                f"sum to {_fmt(rem, i['unit'])}")
    return ""


def _footer_parts(i: dict[str, Any]) -> list[str]:
    """total, share, truncation - the same three sentences on the slide and on the page."""
    parts = [f"total {_fmt(i['total'], i['unit'])} over {i['n_groups']} groups", _share_text(i)]
    rem = _remainder_text(i)
    if rem:
        parts.append(rem)
    return parts


def _title_text(report: dict[str, Any]) -> str:
    return f"{report['database']}: what the data says"


def _provenance_text(report: dict[str, Any]) -> str:
    o = report["ontology"]
    return (f"{o['objects']} objects, {o['links']} measured links, "
            f"{o['measures']} declared measures. Every figure compiled and conserved.")


# --------------------------------------------------------------------------
# validation: refuse a report the renderers cannot carry verbatim
# --------------------------------------------------------------------------


def _is_number(x: Any) -> bool:
    return isinstance(x, (int, float)) and not isinstance(x, bool) and math.isfinite(x)


def _is_count(x: Any) -> bool:
    return isinstance(x, int) and not isinstance(x, bool) and x >= 0


def validate_report(report: Any) -> list[str]:
    """Every way the report cannot be briefed, as '<insight id>: <what is wrong>' lines.

    Empty means the renderers will neither crash nor re-round: every figure is
    a finite number already at 2 dp, every insight has rows, a headline and
    the fields the slide interpolates, and the ontology counts are counts.
    """
    problems: list[str] = []
    if not isinstance(report, dict):
        return ["report: not a JSON object"]
    if not isinstance(report.get("database"), str) or not report["database"].strip():
        problems.append("report: 'database' missing or empty")
    if not isinstance(report.get("scope"), str):
        problems.append("report: 'scope' missing or not a string")
    o = report.get("ontology")
    if not isinstance(o, dict):
        problems.append("report: 'ontology' missing or not an object")
    else:
        for k in ("objects", "links", "measures"):
            if not _is_count(o.get(k)):
                problems.append(f"report: ontology.{k} is not a non-negative integer")
    insights = report.get("insights")
    if not isinstance(insights, list) or not insights:
        problems.append("report: 'insights' missing or empty; nothing to brief")
        insights = []
    for k, i in enumerate(insights):
        label = (i.get("id") if isinstance(i, dict) and isinstance(i.get("id"), str)
                 and i["id"] else f"insight #{k}")
        if not isinstance(i, dict):
            problems.append(f"{label}: not an object")
            continue
        for f in ("id", "headline", "measure", "dimension", "sql"):
            if not isinstance(i.get(f), str) or not i[f].strip():
                problems.append(f"{label}: '{f}' is missing or empty")
        for f in ("kind", "unit"):
            if not isinstance(i.get(f), str):
                problems.append(f"{label}: '{f}' is missing or not a string")
        rows = i.get("rows")
        if not isinstance(rows, list) or not rows:
            problems.append(f"{label}: no rows; an insight with nothing to tabulate is not briefed")
            rows = []
        for r, row in enumerate(rows):
            if not isinstance(row, (list, tuple)) or len(row) != 2:
                problems.append(f"{label}: row {r} is not a [label, value] pair")
            elif not _is_number(row[1]):
                problems.append(f"{label}: row {r} value is not a finite number")
            elif round(row[1], 2) != row[1]:
                problems.append(f"{label}: row {r} value {row[1]!r} is not at 2 dp; "
                                "the brief renders 2 dp verbatim and will not re-round")
        total = i.get("total")
        if not _is_number(total):
            problems.append(f"{label}: 'total' is not a finite number")
        elif round(total, 2) != total:
            problems.append(f"{label}: total {total!r} is not at 2 dp")
        elif rows and all(isinstance(w, (list, tuple)) and len(w) == 2 and _is_number(w[1])
                          for w in rows):
            if math.fsum(w[1] for w in rows) > total + 0.005 * (len(rows) + 1):
                problems.append(f"{label}: rows sum to more than the total")
        n = i.get("n_groups")
        if not _is_count(n) or n < 1:
            problems.append(f"{label}: 'n_groups' is not a positive integer")
        elif n < len(rows):
            problems.append(f"{label}: n_groups {n} is fewer than the {len(rows)} rows")
        share = i.get("top_share")
        if not _is_number(share) or not 0 <= share <= 1:
            problems.append(f"{label}: 'top_share' is not a number in [0, 1]")
        caveats = i.get("caveats", [])
        if not isinstance(caveats, list) or not all(isinstance(c, str) for c in caveats):
            problems.append(f"{label}: 'caveats' is not a list of strings")
    refusals = report.get("refusals", [])
    if not isinstance(refusals, list):
        problems.append("report: 'refusals' is not a list")
    else:
        for k, r in enumerate(refusals):
            if not isinstance(r, dict) or not all(
                isinstance(r.get(f), str) for f in ("question", "reason", "detail")
            ):
                problems.append(f"refusal #{k}: needs string question, reason and detail")
    return problems


# --------------------------------------------------------------------------
# the two renderers
# --------------------------------------------------------------------------


def build_html(report: dict[str, Any]) -> str:
    o = report["ontology"]
    parts = [
        f"<h1>{_esc(report['database'])} - what the data says, and how we know</h1>",
        f"<p><b>Provenance.</b> {_esc(o['objects'])} object types, {_esc(o['links'])} measured "
        f"links, {_esc(o['measures'])} declared measures. {_esc(report['scope'])}.</p>",
    ]
    for i in report["insights"]:
        parts.append(f"<h2>{_esc(i['headline'])}</h2>")
        parts.append(f"<p>{_esc(i['measure'])} by {_esc(i['dimension'])}: "
                     + "; ".join(_esc(p) for p in _footer_parts(i)) + ".</p>")
        parts.append(f"<table><tr><th>{_esc(i['dimension'])}</th>"
                     f"<th>{_esc(i['measure'])} ({_esc(i['unit'])})</th></tr>")
        for lbl, v in _shown(i):
            parts.append(f"<tr><td>{_esc(lbl)}</td><td>{_esc(_num(v))}</td></tr>")
        parts.append("</table>")
        for c in i.get("caveats", []):
            parts.append(f"<p><i>{_esc(c)}</i></p>")
        parts.append(f"<details><summary>SQL</summary><pre>{_esc(i['sql'])}</pre></details>")
    if report.get("refusals"):
        parts.append("<h2>Questions the ontology refused</h2><ul>")
        for r in report["refusals"][:MAX_REFUSALS]:
            parts.append(f"<li><b>{_esc(r['question'])}</b> - {_esc(r['reason'])}: "
                         f"{_esc(r['detail'][:160])}</li>")
        parts.append("</ul>")
    return "\n".join(parts)


def build_pptx(report: dict[str, Any], out: Path) -> int:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]

    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = _title_text(report)
    s.placeholders[1].text = _provenance_text(report)

    n = 0
    for i in report["insights"]:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = str(i["headline"])
        # paragraph-level font: an empty string has no run to style
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True

        shown = _shown(i)
        table = slide.shapes.add_table(
            len(shown) + 1, 2, Inches(0.5), Inches(1.8), Inches(6), Inches(0.4) * (len(shown) + 1)
        ).table
        table.cell(0, 0).text = str(i["dimension"])
        table.cell(0, 1).text = f"{i['measure']} ({i['unit']})"
        for r, (lbl, v) in enumerate(shown, start=1):
            table.cell(r, 0).text = str(lbl)
            table.cell(r, 1).text = _num(v)

        foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(9), Inches(1))
        ff = foot.text_frame
        ff.word_wrap = True
        ff.text = "; ".join(_footer_parts(i)) + ". " + "; ".join(i.get("caveats", []))
        ff.paragraphs[0].font.size = Pt(10)
        slide.notes_slide.notes_text_frame.text = "SQL:\n" + str(i["sql"])
        n += 1

    if report.get("refusals"):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.text = "Questions the ontology refused rather than guessed"
        tf.paragraphs[0].font.size = Pt(22)
        for r in report["refusals"][:MAX_REFUSALS]:
            p = tf.add_paragraph()
            p.text = f"{r['question']} - {r['reason']}"
            p.font.size = Pt(12)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return n


# --------------------------------------------------------------------------
# read-back: every number on the deck is one the report accounts for
# --------------------------------------------------------------------------


def allowed_numbers(report: dict[str, Any]) -> set[str]:
    """Every numeric token the report accounts for, in the rendering the deck uses.

    Figures: every row value and total at 2 dp, n_groups, the share as the
    headline states it, the truncation line's counts and remainder, the
    ontology counts. Text the report itself carries and the deck prints
    verbatim may hold digits too - labels, headlines, caveats, names, refusal
    lines - so their tokens are accounted for as well.
    """
    o = report["ontology"]
    allowed: set[str] = set()
    for x in (report["database"], o["objects"], o["links"], o["measures"]):
        allowed |= _numbers(str(x))
    for i in report["insights"]:
        for lbl, v in i["rows"]:
            allowed |= _numbers(str(lbl)) | _numbers(_num(v))
        allowed |= _numbers(_num(i["total"])) | {str(i["n_groups"])}
        allowed |= _numbers(_share_text(i)) | _numbers(_remainder_text(i))
        for text in (i["headline"], i["measure"], i["dimension"], i["unit"],
                     *i.get("caveats", [])):
            allowed |= _numbers(str(text))
    for r in report.get("refusals", [])[:MAX_REFUSALS]:
        allowed |= _numbers(str(r["question"])) | _numbers(str(r["reason"]))
    return allowed


def deck_numbers(path: Path) -> set[str]:
    """Every numeric token on every slide face of the deck - text frames and table cells."""
    from pptx import Presentation

    seen: set[str] = set()
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                seen |= _numbers(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        seen |= _numbers(cell.text)
    return seen


def verify_deck(report: dict[str, Any], path: Path) -> list[str]:
    """Read the deck back; return every number on it that the report does not account for."""
    return sorted(deck_numbers(path) - allowed_numbers(report))


# --------------------------------------------------------------------------
# entry point
# --------------------------------------------------------------------------


def main(argv: list[str] | None = None) -> int:
    # R-0012: a Windows console defaults to cp1252 and turns every accented
    # label into a replacement glyph. Say what the data says.
    for stream in (sys.stdout, sys.stderr):
        if hasattr(stream, "reconfigure"):
            stream.reconfigure(encoding="utf-8", errors="replace")
    ap = argparse.ArgumentParser(
        description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter
    )
    ap.add_argument("--insights", type=Path, required=True)
    ap.add_argument("--pptx", type=Path)
    ap.add_argument("--html", type=Path)
    args = ap.parse_args(argv)
    if not args.pptx and not args.html:
        ap.error("pass --pptx and/or --html; there is nothing to write otherwise")

    try:
        report = json.loads(args.insights.read_text(encoding="utf-8"))
    except (OSError, ValueError) as e:  # json.JSONDecodeError is a ValueError
        print(f"FAIL report: cannot read {args.insights}: {e}")
        return 1
    problems = validate_report(report)
    if problems:
        for p in problems:
            print(f"FAIL {p}")
        return 1

    if args.pptx:
        made = build_pptx(report, args.pptx)
        print(f"  wrote {args.pptx} ({made} insight slides)")
        if made != len(report["insights"]):
            print(f"FAIL {len(report['insights'])} insights, {made} slides")
            return 1
        stray = verify_deck(report, args.pptx)
        if stray:
            print("FAIL the deck, read back, carries figures the report does not: "
                  + ", ".join(stray))
            return 1
    if args.html:
        args.html.parent.mkdir(parents=True, exist_ok=True)
        args.html.write_text(build_html(report), encoding="utf-8")
        print(f"  wrote {args.html}")
    if not args.pptx:
        print("wrote HTML; no deck built")
        return 0
    print(f"PASS {len(report['insights'])} insights on {len(report['insights'])} slides; "
          "every figure on the deck, read back, is in the report")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

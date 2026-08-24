"""The insight miner must report only what conserves, and the deck must carry only what was mined.

Two ways this layer could lie, and each gets a test that proves the lie is caught:

  * an insight whose grouped total does not equal the raw total is the fan-out
    this whole stack exists to make unreachable. mine() must refuse to report
    it and must fail the run, not drop it quietly;
  * a slide that carries a figure the report does not is a retyped number, the
    one artifact a buyer actually holds. The deck is read back and every number
    on it must be present in the insights JSON it was generated from.

Fixtures are built here (R-0002). The lake is not needed.
"""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

import brief  # noqa: E402
import insights as ins  # noqa: E402
from brief import (  # noqa: E402
    _numbers,
    build_html,
    build_pptx,
    deck_numbers,
    validate_report,
    verify_deck,
)


@pytest.fixture()
def lake(tmp_path: Path):  # noqa: ANN201
    """Orders -> territory, with one territory carrying most of the revenue."""
    import duckdb

    c = duckdb.connect(":memory:")
    (tmp_path / "db").mkdir()
    c.execute(
        "COPY (SELECT * FROM (VALUES "
        "(1, 1, 700.0), (2, 1, 200.0), (3, 2, 50.0), (4, 3, 30.0), (5, 4, 20.0), "
        "(6, NULL, 100.0)"
        ") AS t(order_id, territory_id, amount)) TO '"
        + (tmp_path / "db" / "Sales.Orders.parquet").as_posix() + "' (FORMAT PARQUET)"
    )
    c.execute(
        "COPY (SELECT * FROM (VALUES (1, 'North', 'x'), (2, 'South', 'y'), "
        "(3, 'East', NULL), (4, 'West', NULL)"
        ") AS t(territory_id, name, note)) TO '"
        + (tmp_path / "db" / "Sales.Territory.parquet").as_posix() + "' (FORMAT PARQUET)"
    )
    manifest = {
        "database": "Fixture",
        "tables": [
            {"schema": "Sales", "table": "Orders", "declared_rows": 6, "extracted_rows": 6,
             "columns": 3, "path": "db/Sales.Orders.parquet"},
            {"schema": "Sales", "table": "Territory", "declared_rows": 4, "extracted_rows": 4,
             "columns": 3, "path": "db/Sales.Territory.parquet"},
        ],
        "skipped": [],
        "primary_keys": {"Sales.Orders": ["order_id"], "Sales.Territory": ["territory_id"]},
        "foreign_keys": [
            {"name": "FK_T", "from_table": "Sales.Orders", "from_column": "territory_id",
             "to_table": "Sales.Territory", "to_column": "territory_id"},
        ],
    }
    try:
        yield c, manifest, tmp_path
    finally:
        c.close()


def _with_measure(monkeypatch: pytest.MonkeyPatch, root: Path) -> None:
    monkeypatch.setattr(ins, "ROOT", root)
    monkeypatch.setitem(
        ins.MEASURES, "Fixture",
        [{"name": "revenue", "grain": "Sales.Orders", "expr": 'ROUND(SUM(f."amount"), 2)',
          "unit": "USD", "is": "order amount", "is_not": "not net of anything"}],
    )


def test_concentration_is_reported_with_a_conserving_total(lake, monkeypatch) -> None:  # noqa: ANN001
    c, manifest, root = lake
    _with_measure(monkeypatch, root)
    report = ins.mine(c, manifest, top=5)
    assert "error" not in report
    kinds = {i["kind"] for i in report["insights"]}
    assert "dominance" in kinds or "concentration" in kinds
    for i in report["insights"]:
        assert i["conserves"]
        # rows are the top labelled groups; the unlabelled remainder is reported
        # separately, so rows sum to at most the total, never more
        assert sum(v for _, v in i["rows"]) <= i["total"] + 0.02
        assert 0 < i["top_share"] <= 1.0
        assert i["sql"].strip().upper().startswith("SELECT"), "every insight carries its SQL"


def test_an_unknown_bucket_means_a_used_dimension_missing_on_the_measure(
    lake, monkeypatch,  # noqa: ANN001
) -> None:
    """100 of 1100 lands on an order with no territory: 9.1 pct, reported."""
    c, manifest, root = lake
    _with_measure(monkeypatch, root)
    report = ins.mine(c, manifest, top=10)
    unknown = [i for i in report["insights"] if i["kind"] == "unknown_bucket"]
    assert unknown, "the NULL-territory share must be surfaced as a data gap"
    assert "9.1 pct" in unknown[0]["headline"]


def test_a_mostly_empty_column_is_not_a_dimension(lake, monkeypatch) -> None:  # noqa: ANN001
    """Territory.note is NULL on half the rows: an annotation, not a segmentation.

    The first run of the miner reported that 95 pct of sales landed on rows a
    mostly-NULL Suffix column could not label. That is not an insight.
    """
    c, manifest, root = lake
    _with_measure(monkeypatch, root)
    rel = "read_parquet('" + (root / "db" / "Sales.Territory.parquet").as_posix() + "')"
    assert ins._label_attrs(c, rel) == ["name"]


def test_a_long_text_column_is_not_a_dimension(tmp_path: Path) -> None:
    import duckdb

    c = duckdb.connect(":memory:")
    try:
        c.execute("CREATE TABLE t AS SELECT i AS id, repeat('x', 200) || i::VARCHAR AS blob, "
                  "'A' AS tag FROM range(1, 10) r(i)")
        assert ins._label_attrs(c, "t") == []  # tag has 1 distinct; blob is too long
        c.execute("UPDATE t SET tag = 'B' WHERE id > 5")
        assert ins._label_attrs(c, "t") == ["tag"]
    finally:
        c.close()


def test_a_non_conserving_grouping_fails_the_run_rather_than_being_dropped(
    lake, monkeypatch,  # noqa: ANN001
) -> None:
    """Force a fan-out past the compiler and prove mine() refuses to report it.

    The compiler itself cannot produce this, so the test swaps in a compile that
    returns an inflating join. If mine() ever reports such a grouping, or hides
    it, the control is decoration.
    """
    c, manifest, root = lake
    _with_measure(monkeypatch, root)
    from ontology import CompiledQuery, Ontology

    orders = (root / "db" / "Sales.Orders.parquet").as_posix()
    real_compile = Ontology.compile

    def inflating(self, measure, **kw):  # noqa: ANN001, ANN202
        got = real_compile(self, measure, **kw)
        if isinstance(got, CompiledQuery) and kw.get("group_by"):
            # a self-join that doubles every row
            sql = (f'SELECT d0."name" AS label, ROUND(SUM(f."amount"), 2) AS revenue '
                   f"FROM read_parquet('{orders}') f "
                   f"LEFT JOIN (SELECT t.territory_id, t.name FROM read_parquet('"
                   + (root / "db" / "Sales.Territory.parquet").as_posix()
                   + "') t UNION ALL SELECT t.territory_id, t.name FROM read_parquet('"
                   + (root / "db" / "Sales.Territory.parquet").as_posix()
                   + "') t) d0 ON f.territory_id = d0.territory_id "
                   'GROUP BY d0."name"')
            return CompiledQuery(sql=sql, measure=got.measure, grain=got.grain,
                                 group_by=got.group_by, notes=got.notes)
        return got

    monkeypatch.setattr(Ontology, "compile", inflating)
    report = ins.mine(c, manifest, top=5)
    assert report["broken"], "the inflated grouping must be named as broken"
    assert not [i for i in report["insights"] if i["dimension"].endswith("name")], (
        "an inflated grouping was reported as an insight"
    )


# --------------------------------------------------------------------------
# the deck carries the report's figures and nothing else
# --------------------------------------------------------------------------


def _report(**over) -> dict:  # noqa: ANN003
    """A valid one-insight report, built by hand so the brief is tested without the miner."""
    rows = [["North", 700.0], ["South", 200.0], ["East", 100.0]]
    insight = {
        "id": "Fixture-001", "kind": "dominance",
        "headline": "North alone is 70.0 pct of revenue across 3 territory values",
        "measure": "revenue", "unit": "USD", "grain": "Sales.Orders", "dimension": "territory",
        "rows": rows, "total": 1000.0, "n_groups": 3, "top_share": 0.9, "surprise": 0.37,
        "sql": "SELECT 1", "notes": [], "caveats": ["revenue: order amount; not net of anything"],
        "conserves": True,
    }
    insight.update(over.pop("insight", {}))
    report = {
        "database": "Fixture", "ontology": {"objects": 2, "links": 1, "measures": 1},
        "insights": [insight], "refusals": [], "broken": [],
        "scope": "every figure compiled by the ontology",
    }
    report.update(over)
    return report


def _deck_text(path: Path) -> list[str]:
    """Every text frame and table cell on every slide, in order."""
    from pptx import Presentation

    out: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if shape.has_table:
                out.extend(cell.text for row in shape.table.rows for cell in row.cells)
    return out


def _retype(path: Path, text: str) -> None:
    """Tamper with the first table cell value on the first insight slide - a retyped figure."""
    from pptx import Presentation

    prs = Presentation(str(path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                shape.table.cell(1, 1).text = text
                prs.save(str(path))
                return
    raise AssertionError("no table on the deck")


def test_numbers_matches_integers_and_one_decimal_tokens_not_only_two_decimals() -> None:
    """Finding 24: '12345 groups' and '86.2 pct' were invisible to the old read-back regex."""
    text = ("12345 groups, 86.2 pct, 1,234.56 USD, showing 8 of 35 rows, -0.30 USD, "
            "AdventureWorks2025 AddressLine1 over 35. end")
    assert _numbers(text) == {"12345", "86.2", "1234.56", "8", "35", "0.30"}


def test_every_figure_on_the_deck_is_in_the_report(lake, monkeypatch, tmp_path: Path) -> None:  # noqa: ANN001
    c, manifest, root = lake
    _with_measure(monkeypatch, root)
    report = ins.mine(c, manifest, top=5)
    out = tmp_path / "brief.pptx"
    n = build_pptx(report, out)
    assert n == len(report["insights"])

    stray = verify_deck(report, out)
    assert not stray, f"figures on the deck that are not in the report: {stray}"
    seen = deck_numbers(out)
    assert seen, "the deck carried no figures at all"
    # and the other direction, built here without the brief's helpers: every
    # shown row value, every total and every group count is on the deck
    for i in report["insights"]:
        for _, v in i["rows"][:8]:
            assert f"{v:.2f}" in seen
        assert f"{i['total']:.2f}" in seen
        assert str(i["n_groups"]) in seen
        pct = re.search(r"(\d+\.\d) pct", i["headline"])
        assert pct and pct.group(1) in seen


def test_verify_deck_catches_a_retyped_figure(tmp_path: Path) -> None:
    """The gate must be able to fail (R-0007): retype one cell and it is named."""
    out = tmp_path / "brief.pptx"
    build_pptx(_report(), out)
    assert verify_deck(_report(), out) == []
    _retype(out, "999,999.99")
    assert verify_deck(_report(), out) == ["999999.99"]


def test_rows_beyond_the_cap_are_stated_not_dropped(tmp_path: Path) -> None:
    """Finding 20: the deck showed 8 rows and said only 'over N groups'."""
    rows = [[f"g{k}", float(100 - k)] for k in range(12)]
    total = float(sum(v for _, v in rows))
    top = sum(v for _, v in rows[:3]) / total
    report = _report(insight={
        "kind": "concentration",
        "headline": f"3 of 12 territory values (g0, g1, g2) carry {top * 100:.1f} pct of revenue",
        "rows": rows, "total": total, "n_groups": 12, "top_share": round(top, 4),
    })
    remainder = f"{total - sum(v for _, v in rows[:8]):,.2f}"  # 362.00, computed here
    out = tmp_path / "brief.pptx"
    build_pptx(report, out)
    texts = _deck_text(out)
    footer = next(t for t in texts if t.startswith("total "))
    assert "showing 8 of 12 groups" in footer
    assert (f"the other 4 groups and any rows with no territory value sum to {remainder} USD"
            in footer)
    assert sum(1 for t in texts if t.startswith("g")) == 8, "the table shows exactly 8 labels"
    assert verify_deck(report, out) == []
    page = build_html(report)
    assert "showing 8 of 12 groups" in page and f"sum to {remainder} USD" in page
    assert page.count("<tr><td>g") == 8
    # all rows shown: no truncation line at all
    page = build_html(_report())
    assert "showing" not in page and "sum to" not in page


def test_one_formatter_renders_two_decimals_on_the_page_and_the_deck(tmp_path: Path) -> None:
    """Findings 21/27: the HTML rounded to 0 dp ('USD 2' for 2.5) while the deck used 2 dp."""
    cases = [2.5, -0.3, 43909437.51]
    for v in cases:
        report = _report(insight={"rows": [["North", v]], "total": v, "n_groups": 1})
        page = build_html(report)
        out = tmp_path / "brief.pptx"
        build_pptx(report, out)
        cell = [t for t in _deck_text(out) if t.replace(",", "").replace("-", "")
                .replace(".", "").isdigit()]
        expected = f"{v:,.2f}"
        assert f"total {expected} USD" in page, page
        assert f"<td>{expected}</td>" in page
        assert cell == [expected], cell
        assert f"USD {v:,.0f}" not in page, "the 0 dp rendering is gone"  # USD 2, USD -0


def test_html_escapes_every_interpolated_field() -> None:
    """Finding 22: labels, counts, names and dimensions went into the page raw."""
    report = _report(
        ontology={"objects": "<img src=x onerror=alert(1)>", "links": 1, "measures": 1},
        scope="<svg onload=alert(1)>",
        insight={
            "rows": [["<script>alert(1)</script>", 700.0], ["South", 300.0]],
            "total": 1000.0, "n_groups": 2, "measure": "<i>revenue</i>",
            "dimension": "<b>territory</b>", "unit": "<u>USD</u>",
            "headline": "<script>x</script> alone is 70.0 pct of revenue across 2 values",
            "caveats": ["<marquee>caveat</marquee>"], "sql": "SELECT '<x>'",
        },
        refusals=[{"question": "<q>", "reason": "<r>", "detail": "<d>"}],
    )
    page = build_html(report)
    for raw in ("<script>", "<img", "<svg", "<b>territory", "<i>revenue", "<u>USD", "<marquee>",
                "<q>", "<r>", "<d>", "<x>"):
        assert raw not in page, raw
    for escaped in ("&lt;script&gt;", "&lt;img", "&lt;svg", "&lt;b&gt;territory",
                    "&lt;marquee&gt;", "&lt;q&gt;", "&lt;x&gt;"):
        assert escaped in page, escaped


def test_the_footer_share_is_the_headline_share_verbatim(tmp_path: Path) -> None:
    """Finding 23: the headline said 50.1 pct, the footer re-rendered top_share as 50.0."""
    report = _report(insight={
        "kind": "concentration",
        "headline": "2 of 4 territory values (North, South) carry 50.1 pct of revenue",
        "rows": [["North", 300.0], ["South", 200.5], ["East", 250.0], ["West", 249.5]],
        "total": 1000.0, "n_groups": 4, "top_share": 0.5005,
    })
    assert f"{report['insights'][0]['top_share'] * 100:.1f}" == "50.0", "the fixture disagrees"
    out = tmp_path / "brief.pptx"
    build_pptx(report, out)
    footer = next(t for t in _deck_text(out) if t.startswith("total "))
    assert "top share 50.1 pct" in footer and "50.0 pct" not in footer
    page = build_html(report)
    assert "top share 50.1 pct" in page and "50.0 pct" not in page
    assert verify_deck(report, out) == []


def test_a_broken_report_fails_with_the_insight_named_not_a_traceback(
    tmp_path: Path, capsys: pytest.CaptureFixture[str],
) -> None:
    """Findings 25/26a: an empty headline, a None figure or no rows must be refused by name."""
    cases = {
        "'headline' is missing or empty": {"headline": ""},
        "'total' is not a finite number": {"total": None},
        "no rows": {"rows": []},
        "value 1.005 is not at 2 dp": {"rows": [["North", 1.005]], "total": 1.01, "n_groups": 1},
        "rows sum to more than the total": {"total": 900.0},
        "n_groups 1 is fewer than the 3 rows": {"n_groups": 1},
    }
    for expect, over in cases.items():
        report = _report(insight=over)
        assert any(expect in p for p in validate_report(report)), (expect, validate_report(report))
        src = tmp_path / "r.json"
        src.write_text(json.dumps(report), encoding="utf-8")
        rc = brief.main(["--insights", str(src), "--pptx", str(tmp_path / "d.pptx")])
        out = capsys.readouterr().out
        assert rc == 1
        assert "FAIL Fixture-001:" in out and expect in out, out
        assert "PASS" not in out and "Traceback" not in out
    assert "report: 'ontology' missing" in " ".join(validate_report(_report(ontology=None)))
    assert validate_report(_report()) == []


def test_html_only_writes_the_page_and_does_not_claim_pass(
    tmp_path: Path, capsys: pytest.CaptureFixture[str],
) -> None:
    """Finding 26c: with --html alone nothing is read back, so nothing may be called PASS."""
    src = tmp_path / "r.json"
    src.write_text(json.dumps(_report()), encoding="utf-8")
    page = tmp_path / "b.html"
    assert brief.main(["--insights", str(src), "--html", str(page)]) == 0
    out = capsys.readouterr().out
    assert "wrote HTML; no deck built" in out
    assert "PASS" not in out
    assert "what the data says, and how we know" in page.read_text(encoding="utf-8")


def test_pass_is_printed_only_after_the_deck_reads_back_clean(
    tmp_path: Path, capsys: pytest.CaptureFixture[str], monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Finding 26b/d: main() reads the deck back; a retyped figure turns PASS into FAIL."""
    src = tmp_path / "r.json"
    src.write_text(json.dumps(_report()), encoding="utf-8")
    deck = tmp_path / "d.pptx"
    assert brief.main(["--insights", str(src), "--pptx", str(deck), "--html",
                       str(tmp_path / "b.html")]) == 0
    out = capsys.readouterr().out
    assert out.rstrip().splitlines()[-1].startswith("PASS")

    real = brief.build_pptx

    def tampered(report, path):  # noqa: ANN001, ANN202
        n = real(report, path)
        _retype(path, "701.00")
        return n

    monkeypatch.setattr(brief, "build_pptx", tampered)
    assert brief.main(["--insights", str(src), "--pptx", str(deck)]) == 1
    out = capsys.readouterr().out
    assert "FAIL" in out and "701.00" in out and "PASS" not in out


def test_the_html_brief_carries_the_sql_behind_each_figure(lake, monkeypatch) -> None:  # noqa: ANN001
    c, manifest, root = lake
    _with_measure(monkeypatch, root)
    report = ins.mine(c, manifest, top=5)
    page = build_html(report)
    for i in report["insights"]:
        assert i["headline"] in page or i["headline"].replace("'", "&#x27;") in page
        assert "<pre>" in page and "SELECT" in page
    assert "what the data says, and how we know" in page
